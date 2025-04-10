from math import floor
import random
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import pandas as pd
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
import requests
from io import BytesIO


class CustomPresentation:
    def __init__(self, key, name):
        self.prs = Presentation()
        self.key = key
        self.name = name

        self._configure_presentation()

    def _configure_presentation(self):
        """Configures the presentation dimensions and layout."""
        self.prs.slide_width = Inches(15)
        self.prs.slide_height = Inches(8.5)
        self.slide_layout = self.prs.slide_layouts[5]  # Blank layout

        # Calculate dimensions
        self.slide_width = self.prs.slide_width - 10
        self.slide_height = self.prs.slide_height - 10
        self.col_width = self.slide_width / 4
        self.img_width = (self.col_width * 3) + Inches(1)
        self.img_height = (self.slide_height / 2) - Inches(0.5)

        # Ribbon
        self.ribbon_color = RGBColor(115, 204, 255)

        # Comment Boxes
        self.issue_box_color = RGBColor(0, 102, 161)
        self.because_box_color = RGBColor(203, 211, 224)

        # Logo Paths
        self.logo = "assets/logo.png"

        self._add_intro_slide()

    def _add_intro_slide(self):
        """Adds an introduction slide."""
        slide = self._get_new_slide()

        # Apply a gradient background from top left (dark blue) to bottom right (turquoise)
        self._add_gradient(slide)

        intro_texts = [
            ("Player Centuries Scored Analysis", 48, Inches(3), False),
            (
                f"{self.name} in Professional Cricket",
                32,
                Inches(3.6),
                False,
            ),
            ("A PPT Automation Demo", 24, Inches(4.5), False),
            ("By: Ayush - Automation", 16, Inches(4.85), False),
            (
                "dhanrajaayush123@gmail.com",
                16,
                Inches(5.08),
                False,
            ),
        ]

        left = Inches(1)

        for text, font_size, top, bold in intro_texts:
            text_box = slide.shapes.add_textbox(left, top, self.slide_width, Inches(1))
            text_frame = text_box.text_frame
            text_frame.text = text
            p = text_frame.paragraphs[0]
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.name = "Calibri Light"
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.LEFT

        #  innovation logo
        top = Inches(6)
        left = Inches(1)
        picture = slide.shapes.add_picture(self.logo, left, top)
        picture.width = int(picture.width / 3)
        picture.height = int(picture.height / 3)

        #  logo
        self._add_slide_logo(slide, self.logo)

    def add_end_slide(self):
        """Adds End slide."""
        slide = self._get_new_slide()
        self._add_gradient(slide)

        # Load the white logo image
        logo_path = self.logo

        # Calculate the size and position for the logo
        logo_picture = slide.shapes.add_picture(logo_path, 0, 0)

        # Center the logo horizontally and vertically
        logo_left = int((self.prs.slide_width - logo_picture.width) / 2)
        logo_top = int((self.prs.slide_height - logo_picture.height) / 2)

        # Set the logo position
        logo_picture.left = logo_left
        logo_picture.top = logo_top

    def _add_gradient(self, slide):
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 325

        # Set gradient colors
        stop_1 = fill.gradient_stops[0]
        stop_1.position = 0  # Position at the start (top left)
        stop_1.color.rgb = RGBColor(0, 52, 120)  # Dark blue

        stop_2 = fill.gradient_stops[1]
        stop_2.position = 1  # Position at the end (bottom right)
        stop_2.color.rgb = RGBColor(0, 137, 196)  # Turquoise

    def _get_new_slide(self):
        slide = self.prs.slides.add_slide(self.slide_layout)

        if slide.shapes.title:
            # Delete the title placeholder
            sp = slide.shapes.title
            slide.shapes._spTree.remove(sp._element)

        return slide

    def SubElement(self, parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def _set_shape_transparency(self, shape, alpha):
        """Set the transparency (alpha) of a shape"""
        ts = shape.fill._xPr.solidFill
        sF = ts.get_or_change_to_srgbClr()
        sE = self.SubElement(sF, "a:alpha", val=str(alpha))

    def wrap_text(self, text, max_length):
        """Wrap text into a maximum of two lines if it exceeds the max_length."""
        lines = []
        line_count = 0

        while len(text) > max_length and line_count < 2:
            # Find the break point for the line
            break_point = max_length
            # Adjust break_point to not split words in the middle
            while break_point > 0 and text[break_point] not in " \n":
                break_point -= 1
            if break_point == 0:  # If no space found, split at max_length
                break_point = max_length
            lines.append(text[:break_point].strip())
            text = text[break_point:].strip()
            line_count += 1

        if text:
            if line_count < 2:
                lines.append(text)
            else:
                # If the text exceeds two lines, append the remaining text to the second line
                lines[-1] += "..."

        return lines

    def preprocess_data(self, data, max_length=30):
        """Preprocess data to wrap text in each cell, ensure all columns have the same length,
        and remove purely numeric entries only from specific columns."""
        wrapped_data = {key: [] for key in data}
        max_rows = 0

        # Process each column separately
        for key in data:
            for item in data[key]:
                if isinstance(item, str):
                    item = item.strip()  # Remove leading and trailing whitespace
                    if key == "Product Descriptions":
                        # Remove purely numeric entries from 'Product Descriptions'
                        if not item.replace(".", "", 1).isdigit():
                            wrapped_lines = self.wrap_text(item, max_length)
                            wrapped_data[key].extend(wrapped_lines)
                    else:
                        wrapped_data[key].append(item)
                else:
                    wrapped_data[key].append(str(item))
            max_rows = max(max_rows, len(wrapped_data[key]))

        # Ensure all columns have the same length
        for key in wrapped_data:
            while len(wrapped_data[key]) < max_rows:
                wrapped_data[key].append("")

        return wrapped_data

    def _style_table(self, table):
        """Applies styling to a table."""
        table.auto_set_font_size(False)
        table.set_fontsize(10)  # Set font size
        table.scale(1.2, 1.2)  # Scale table to fit size

        for key in table.get_celld().keys():
            cell = table.get_celld()[key]
            if key[0] == 0:  # Header row
                cell.set_text_props(weight="bold", color="white")
                cell.set_facecolor("#73CCFF")  # Header background color
            else:  # Data rows
                cell.set_facecolor("#F2F2F2")  # Row background color
            cell.set_edgecolor(cell.get_facecolor())

    def _add_top_text_boxes(
        self,
        slide,
        slide_name,
        country=None,
    ):
        """Adds text boxes at the top of the slide with borders, background color, and text wrapping.

        Handles both the case with a single box spanning the entire width and four boxes splitting the width."""
        box_height = Inches(0.5)
        top = Inches(0)

        if country is None:
            # Single box spanning the entire width
            width = self.prs.slide_width
            self._create_text_box(
                slide, 0, top, width, box_height, slide_name, PP_ALIGN.LEFT, True
            )
        else:
            # Four boxes
            width = self.prs.slide_width / 2

            self._create_text_box(
                slide, 0, top, width, box_height, slide_name, PP_ALIGN.CENTER
            )
            self._create_text_box(
                slide,
                width,
                top,
                width,
                box_height,
                ", ".join(country),
                PP_ALIGN.CENTER,
            )

    def _create_text_box(
        self, slide, left, top, width, height, text, align=PP_ALIGN.CENTER, long=False
    ):
        """Helper function to create and style a text box with dynamic font size adjustment."""
        # Create the main text box
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # Initial font size
        font_size = Pt(16)

        para = text_frame.paragraphs[0]

        para.text = text
        para.font.size = font_size
        para.alignment = align

        # Set border
        shape = text_box.line
        shape.color.rgb = RGBColor(255, 255, 255)
        shape.width = Pt(2)

        # Set background color
        fill = text_box.fill
        fill.solid()
        fill.fore_color.rgb = self.ribbon_color

    def add_player_image_from_url(
        self,
        slide,
        image_url,
        left=Inches(2),
        top=Inches(1),
        width=Inches(4),
        height=Inches(5),
        border_color=RGBColor(0, 102, 204),
        border_width=Pt(3),
    ):
        """
        Downloads an image from a URL and places it on the left with cover-style scaling and a border.
        """
        try:
            response = requests.get(image_url)
            if response.status_code == 200:
                image_stream = BytesIO(response.content)

                # Add a border frame
                frame = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, left, top, width, height
                )
                frame.fill.solid()
                frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
                frame.fill.transparency = 1
                frame.line.color.rgb = border_color
                frame.line.width = border_width

                # Add the image on top (same position and size)
                slide.shapes.add_picture(
                    image_stream, left, top, width=width, height=height
                )
            else:
                print(f"Failed to download image from {image_url}")
        except Exception as e:
            print(f"Error downloading image: {e}")

    def add_player_info(self, slide_name, personal_info, total_cen):
        country = personal_info["Country"].unique()
        image = personal_info["Image"].unique()[0]
        flag = personal_info["Flag"].unique()[0]

        slide = self._get_new_slide()
        boundary = (Inches(-2), Inches(8), Inches(15), Inches(9))
        self._add_random_boxes(slide, boundary, 100)
        self._add_top_text_boxes(slide, slide_name, country)
        self.add_player_image_from_url(slide, image)
        self.add_player_image_from_url(
            slide,
            flag,
            left=Inches(4.8),
            top=Inches(5.8),
            width=Inches(2),
            height=Inches(1),
        )

        dob_data = pd.DataFrame(
            [
                {
                    "DOB": personal_info["DOB"].unique()[0]
                    if len(personal_info["DOB"].unique())
                    else "Unknown"
                }
            ]
        )

        self._add_score_table(
            slide,
            dob_data,
            Inches(2),
            top=Inches(1),
            right=Inches(5),
        )

        birth_place = pd.DataFrame(
            [
                {
                    "Birth Place": personal_info["Birth Place"].unique()[0]
                    if len(personal_info["Birth Place"].unique())
                    else "Unknown"
                }
            ]
        )

        self._add_score_table(
            slide,
            birth_place,
            Inches(2),
            top=Inches(1),
            right=Inches(2),
        )

        mother = pd.DataFrame(
            [
                {
                    "Mother": personal_info["Mother"].unique()[0]
                    if len(personal_info["Mother"].unique())
                    else "No Data"
                }
            ]
        )

        self._add_score_table(
            slide,
            mother,
            Inches(2),
            top=Inches(2),
            right=Inches(5),
        )

        father = pd.DataFrame(
            [
                {
                    "Father": personal_info["Father"].unique()[0]
                    if len(personal_info["Father"].unique())
                    else "No Data"
                }
            ]
        )

        self._add_score_table(
            slide,
            father,
            Inches(2),
            top=Inches(2),
            right=Inches(2),
        )

        height = pd.DataFrame(
            [
                {
                    "Height": personal_info["Height"].unique()[0]
                    if len(personal_info["Height"].unique())
                    else "No Data"
                }
            ]
        )

        self._add_score_table(
            slide,
            height,
            Inches(2),
            top=Inches(3),
            right=Inches(5),
        )

        married = pd.DataFrame(
            [
                {
                    "Marital Status": personal_info["Marital Status"].unique()[0]
                    if len(personal_info["Marital Status"].unique())
                    else "No Data"
                }
            ]
        )

        self._add_score_table(
            slide,
            married,
            Inches(2),
            top=Inches(3),
            right=Inches(2),
        )

        retired = pd.DataFrame(
            [
                {
                    "Retired": personal_info["Retired"].unique()[0]
                    if len(personal_info["Retired"].unique())
                    else "No Data"
                }
            ]
        )

        self._add_score_table(
            slide,
            retired,
            Inches(2),
            top=Inches(4),
            right=Inches(3.5),
        )

        total_cen = pd.DataFrame([{"Total Centuries": total_cen}])

        self._add_score_table(
            slide,
            total_cen,
            Inches(2),
            top=Inches(5),
            right=Inches(3.5),
        )

        self._add_slide_logo(slide, self.logo)

    def add_slide(
        self,
        slide_name,
        img_paths,
        player_df,
        html_graph_filename="",
    ):
        country = player_df["country"].unique()
        score_data = (
            player_df.groupby("date", as_index=False)["Score"]
            .count()
            .rename(columns={"date": "Year", "Score": "Centuries"})
            .sort_values(by="Year", ascending=True)
            .reset_index(drop=True)
        )
        max_score_data = player_df[player_df["Score"] == player_df["Score"].max()][
            ["date", "Score"]
        ].rename(columns={"date": "Top Year", "Score": "Top Run"})

        """Adds a new slide with images and optional table."""
        # Check if a table slide needs to be added

        # Add the main slide with images
        slide = self._get_new_slide()

        self._add_top_text_boxes(slide, slide_name, country)

        self._add_slide_logo(slide, self.logo)
        self._add_images(slide, img_paths, html_graph_filename)
        prev_height = self._add_score_table(slide, score_data, Inches(2), top=Inches(1))
        self._add_score_table(
            slide,
            max_score_data,
            Inches(2),
            top=Inches(prev_height + 1),
        )

    def _add_score_table(
        self, slide, score_data, width, top=Inches(1.25), right=Inches(0.3)
    ):
        table_width = width

        left = self.prs.slide_width - right - table_width

        rows, cols = score_data.shape
        table_height = 0.3 * (rows + 1)
        shape = slide.shapes.add_table(
            rows + 1, cols, left, top, table_width, Inches(0.2 * (rows + 1))
        )

        tbl = shape._element.graphic.graphicData.tbl
        style_id = "{69012ECD-51FC-41F1-AA8D-1B2483CD663E}"
        tbl[0][-1].text = style_id

        table = shape.table

        for col_idx, column_name in enumerate(score_data.columns):
            cell = table.cell(0, col_idx)
            cell.text = column_name
            cell.text_frame.paragraphs[0].font.size = Pt(12)

        for row_idx, (index, row) in enumerate(score_data.iterrows(), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(11)

        return table_height

    def _add_slide_logo(self, slide, logoPath, on_top=False):
        """Adds a logo to the top right corner of the slide"""
        top = Inches(0.2) if not on_top else Inches(0)
        right = Inches(0.3)

        picture = slide.shapes.add_picture(logoPath, 0, top)

        picture.width = int(picture.width / 4)
        picture.height = int(picture.height / 4)

        left = self.prs.slide_width - right - picture.width

        picture.left = left

    def _add_images(self, slide, img_paths, html_graph_filename):
        """Adds images to a slide."""
        left = Inches(0.25)
        top = Inches(1)
        for img_path in img_paths:
            slide.shapes.add_picture(
                img_path, left, top, width=self.img_width, height=self.img_height
            )

            # Add a transparent shape on top of the image to serve as the clickable link
            hyperlink_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, left, top, self.img_width, self.img_height
            )

            blueBoxFill = hyperlink_shape.fill
            blueBoxFill.solid()
            blueBoxFillColour = blueBoxFill.fore_color
            blueBoxFillColour.rgb = RGBColor(255, 255, 255)
            self._set_shape_transparency(hyperlink_shape, 1000)
            hyperlink_shape.line.color.rgb = RGBColor(255, 255, 255)

            # Create a hyperlink to the html graph file
            hyperlink_shape.click_action.hyperlink.address = html_graph_filename

            # Position next image below the first
            top += self.img_height

    def _add_random_boxes(self, slide, boundary, n=10):
        """
        Adds n random boxes within the specified boundary to the slide.
        Each box has a random shape, rotation, and fill color.

        Parameters:
        - slide: The slide to which the boxes will be added.
        - boundary: A tuple (x_min, y_min, x_max, y_max) specifying the boundary.
        - n: The number of random boxes to add. Default is 10.
        """
        x_min, y_min, x_max, y_max = boundary

        for _ in range(n):
            # Randomly select a position within the boundary
            left = random.uniform(x_min, x_max)
            top = random.uniform(y_min, y_max)

            # Randomly select a width and height
            # width = Inches(2)  # Inches
            width = Inches(random.uniform(0.5, 2))  # Inches

            # Ensure the box stays within the boundary
            left = min(left, x_max - width)
            top = min(top, y_max - width)

            text_box = slide.shapes.add_textbox(left, top, width, width)

            # Randomly rotate the textbox
            rotation_angle = random.uniform(
                0, 360
            )  # Random rotation between 0 and 360 degrees
            text_box.rotation = rotation_angle

            fill = text_box.fill
            fill.solid()
            randomBlue = random.randint(200, 255)
            randomRed = random.randint(100, 115)
            fill.fore_color.rgb = RGBColor(randomRed, 204, randomBlue)

            transparency = random.randint(0, 100)
            fill.transparency = transparency

    def _add_top_drivers(self, driver_data, key):
        # Step 1: Group driver_data into groups of four
        grouped_data = [driver_data[i : i + 4] for i in range(0, len(driver_data), 4)]

        # Step 2: Iterate through each group and create a slide

        for group in grouped_data:
            slide = self._get_new_slide()  # Create a new slide
            left = Inches(0.5)
            top = Inches(0.1)

            text_box1 = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.5))
            text_frame1 = text_box1.text_frame
            text_frame1.text = f"{key} Calls - Major Warranty Drivers"
            text_frame1.paragraphs[0].font.size = Pt(18)

            boundary = (Inches(-2), Inches(8), Inches(15), Inches(9))
            self._add_random_boxes(slide, boundary, 100)

            # Positioning variables for tables
            top = Inches(2)
            right = Inches(0.7)
            table_width = Inches(3)  # Adjust width if needed

            # Step 3: Add a table for each product in the group
            for product_dict in group:
                for product_name, series_data in product_dict.items():
                    # Convert series to DataFrame
                    score_data = series_data.reset_index()
                    score_data.columns = [product_name, "Count"]

                    # Add the table to the slide
                    self._add_score_table(
                        slide, score_data, width=table_width, top=top, right=right
                    )

                    # Increment top to position the next table below the previous one
                    right += Inches(0.5) + table_width  # Adjust spacing as needed

    def add_aggregate_slide(
        self,
        key,
        table_data,
        table_width=Inches(7.5),
        table_left=Inches(3.5),
        extra_header=None,
        merge_indices=None,
    ):
        slide = self._get_new_slide()

        # Add random Boxes at the bottom
        # x_min, y_min, x_max, y_max
        boundary = (Inches(-2), Inches(8), Inches(15), Inches(9))
        self._add_random_boxes(slide, boundary, 100)

        # Text Frames
        left = Inches(0.5)
        top = Inches(0.1)

        text_box1 = slide.shapes.add_textbox(left, top, Inches(2.5), Inches(0.5))
        text_frame1 = text_box1.text_frame
        text_frame1.text = f"{key} - Player Centuries Scored Analysis"
        text_frame1.paragraphs[0].font.size = Pt(12)

        # Table
        top = Inches(1)
        left = table_left

        rows, cols = table_data.shape
        if extra_header and merge_indices:
            shape = slide.shapes.add_table(
                rows + 2, cols, left, top, table_width, Inches(0.2 * (rows + 2))
            )
        else:
            shape = slide.shapes.add_table(
                rows + 1, cols, left, top, table_width, Inches(0.2 * (rows + 1))
            )

        tbl = shape._element.graphic.graphicData.tbl
        style_id = "{69012ECD-51FC-41F1-AA8D-1B2483CD663E}"
        tbl[0][-1].text = style_id

        table = shape.table

        row_additive = 0

        if extra_header and merge_indices:
            row_additive = 1
            for indexes, text in zip(merge_indices, extra_header):
                print(indexes)
                print(text)
                (start_col, end_col) = indexes
                # Merge cells in the specified range
                cell_start = table.cell(0, start_col)
                cell_end = table.cell(0, end_col - 1)
                cell_start.merge(cell_end)
                # Set the text for the merged cell
                cell_start.text = text
                cell_start.text_frame.paragraphs[0].font.size = Pt(16)
                cell_start.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

        for col_idx, column_name in enumerate(table_data.columns):
            cell = table.cell(row_additive, col_idx)
            cell.text = column_name
            cell.text_frame.paragraphs[0].font.size = Pt(16)

            if col_idx > 0:
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        for row_idx, (index, row) in enumerate(table_data.iterrows(), start=1):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + row_additive, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(14)

                # Align all columns except the first one to the right
                if col_idx > 0:
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

        # Logo
        self._add_slide_logo(slide, self.logo)

        # Re-position slide
        xml_slides = self.prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])  # Remove the last slide (the one just added)
        xml_slides.insert(1, slides[-1])  # Insert it as the 2nd slide (index 1)

    def save(self, filename):
        """Saves the presentation."""
        self.add_end_slide()
        self.prs.save(filename)
